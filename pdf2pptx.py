import argparse
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt

# Function to calculate new slide dimensions while preserving the
# image's aspect ratio
def calculate_slide_size(img_width, img_height):
    # Standard slide size in inches (can be set as any initial value)
    base_width = 10.0  # for example, 10 inches as the base width
    img_ratio = img_width / img_height
    slide_height = base_width / img_ratio
    return Inches(base_width), Inches(slide_height)

# Function to copy notes from one slide to another preserving the
# formatting
def copy_notes(src_slide, dest_slide):
    if not src_slide.notes_slide or not src_slide.notes_slide.notes_text_frame:
        return

    # Clear existing notes in the destination slide
    dest_notes_slide = dest_slide.notes_slide
    dest_text_frame = dest_notes_slide.notes_text_frame
    dest_text_frame.clear()

    # Flag to track if we are on the first paragraph
    first_paragraph = True

    # Copy paragraphs and their formatting
    for paragraph in src_slide.notes_slide.notes_text_frame.paragraphs:
        # Use the existing paragraph if it's the first, otherwise add
        # a new one
        if first_paragraph:
            dest_paragraph = dest_text_frame.paragraphs[0]
            first_paragraph = False
        else:
            dest_paragraph = dest_text_frame.add_paragraph()

        # Copy runs (this preserves formatting like bold, italic, etc.)
        for run in paragraph.runs:
            dest_run = dest_paragraph.add_run()
            dest_run.text = run.text
            dest_run.font.bold = run.font.bold
            dest_run.font.italic = run.font.italic
            dest_run.font.underline = run.font.underline
            dest_run.font.size = run.font.size

            # Check if the color exists and has an rgb attribute
            # before setting it
            if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                dest_run.font.color.rgb = run.font.color.rgb

# Main function
def pdf_to_pptx(pdf_file, pptx_file, skip_first, skip_pages, dpi, notes_pptx=None):
    # Convert the PDF to images
    pages = convert_from_path(pdf_file, dpi=dpi)

    # Skip the first page if the option is specified
    if skip_first:
        pages = pages[1:]

    # Convert the page numbers to skip into indices. No risk of 
    # offset if --skip-first is given because the two options 
    # are incompatible and cannot be given simultaneously.
    skip_indices = set(page_num - 1 for page_num in skip_pages)

    # Create a PowerPoint presentation
    prs = Presentation()

    # Open the <notes>.pptx file if provided
    if notes_pptx:
        notes_prs = Presentation(notes_pptx)
    else:
        notes_prs = None

    # Add each image as a slide, skipping the specified pages
    for i, page in enumerate(pages):
        if i in skip_indices:
            continue

        # Save each page as an image (to access its size)
        img_path = 'temp_image.png'
        page.save(img_path, 'PNG')

        # Get image size (in pixels)
        img_width, img_height = page.size

        
        # Calculate new slide size based on the image's aspect ratio
        slide_width, slide_height = calculate_slide_size(img_width, img_height)

        # Adjust the slide size to match the image's aspect ratio
        prs.slide_width = slide_width
        prs.slide_height = slide_height

        # Add slide and insert image
        slide_layout = prs.slide_layouts[6]  # Blank slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(img_path,
                                 Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)

        # Copy notes from the corresponding slide in <notes>.pptx
        if notes_prs and i < len(notes_prs.slides):
            copy_notes(notes_prs.slides[i], slide)

    # Save the PowerPoint presentation
    prs.save(pptx_file)
    print(f"PPTX presentation created: {pptx_file}")

# Argument parser configuration
if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert a PDF to a PowerPoint presentation")
    
    # Add a mutually exclusive group for --skip-first and --skip
    group = parser.add_mutually_exclusive_group()
    group.add_argument("--skip-first",
                       action="store_true",
                       help="Skip the first page of the PDF")
    group.add_argument("--skip",
                       type=str,
                       help="List of pages to skip, separated by commas (e.g., 2,4,5)")

    parser.add_argument("pdf_file", help="Input PDF file name")
    parser.add_argument("pptx_file", help="Output PPTX file name")
    parser.add_argument("--notes-pptx",
                        type=str, help="Existing PPTX file with notes to copy",
                        required=False)
    parser.add_argument("--dpi",
                        type=int, default=200,
                        help="Resolution for the generated images (default: 200 DPI)")

    args = parser.parse_args()

    # Convert the --skip argument to a list of integers
    if args.skip:
        skip_pages = list(map(int, args.skip.split(',')))
    else:
        skip_pages = []

    # Call the function with the arguments
    pdf_to_pptx(args.pdf_file, args.pptx_file, args.skip_first, skip_pages, args.dpi, args.notes_pptx)
